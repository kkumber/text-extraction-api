import pptx
import pptx.enum.shapes
import io
from typing import List

from services.image_ocr import extract_text_from_image


def extract_text_from_pptx(file: bytes) -> List[str]:
    """Extract text from PPTX file including images and tables."""
    ppt = pptx.Presentation(io.BytesIO(file))
    fullText = []
    
    for slide_num, slide in enumerate(ppt.slides):
        for shape in slide.shapes:
            try:
                # Extract text from shapes
                if hasattr(shape, 'text') and shape.text.strip():
                    fullText.append(shape.text)
                
                # Extract text from images
                elif hasattr(shape, 'shape_type') and shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                    if hasattr(shape, 'image'):
                        image_bytes = shape.image.blob
                        ocr_text = extract_text_from_image(image_bytes)
                        if ocr_text.strip():
                            fullText.append(ocr_text)
                
                # Extract text from tables
                elif hasattr(shape, 'has_table') and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                fullText.append(cell.text)
            except Exception as e:
                print(f"Error processing shape in slide {slide_num}: {e}")
                continue
    
    return fullText