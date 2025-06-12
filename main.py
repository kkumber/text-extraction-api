from fastapi import FastAPI, File, UploadFile, HTTPException, BaseModel
import magic
import pymupdf
import docx
import pytesseract as ocr
import PIL
import io
import pptx

app = FastAPI()

@app.post("/upload-document/")
async def upload_document(file: list[UploadFile] = File(...)):

    # Check if empty upload
    if not file:
        raise HTTPException(400, "No files provided")

    MAX_FILE_SIZE = 10 * 1024 * 1024 # 10MB       

    if len(file) > 10:
        raise HTTPException(400, "Too many files")

    # Result Dictinoary
    result = {
        'results': [],
        'total_files': len(file),
    }
    
    # Dictionary parsers
    allowed_docs = {
        'application/pdf': extract_text_from_pdf,
        'application/msword': extract_text_from_docx,
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': extract_text_from_docx,
        'image/jpeg': extract_text_from_image,
        'image/png': extract_text_from_image,
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': extract_text_from_pptx
    }
    
    for i, uploadedFile in enumerate(file):
        # Check if uploaded file exceeds max file size
        if uploadedFile.size > MAX_FILE_SIZE:
            raise HTTPException(400, f"File {uploadedFile.filename} too large")


        content = await uploadedFile.read()
        mime_type = magic.from_buffer(content, mime=True)
        currentProcessingFile = result['results']
        
        # Check if file type is not allowed
        if mime_type not in allowed_docs:
            raise HTTPException(
                status_code=400,
                detail=f"Document type not allowed. Received: {mime_type}"
            )
        
        try:
            extractedText = allowed_docs[mime_type](content)
            currentProcessingFile.append(
            {
                'filename': uploadedFile.filename,
                'mime_type': mime_type,
                'status': 'success',
                'extractedText': extractedText 
            })

        except Exception as e:
            currentProcessingFile.append({
                'filename': uploadedFile.filename,
                'mime_type': mime_type,
                'status': 'error',
                'error': str(e)
            })
                
    return result

# Extract text in pdf file using pymupdf
def extract_text_from_pdf(file):
    doc = pymupdf.open(stream=file, filetype='pdf')
    fullText = []

    try:
        for page in doc:
            # Get texts
            fullText.append(page.get_text())

            # Process all images on the page at once
            image_list = page.get_images()
            for img in image_list:
                xref = img[0]
                img_data = doc.extract_image(xref)
                img_bytes = img_data["image"]
                fullText.append(extract_text_from_image(img_bytes))

        return fullText
    finally:
        doc.close()

# Extract text in docx file using python-docx
def extract_text_from_docx(file):
    doc = docx.Document(io.BytesIO(file))
    fullText = []

    # Get texts on each paragraph block
    for para in doc.paragraphs:
        fullText.append(para.text)

    # Get image    
    for image in doc.part.related_parts.values():
        if image.content_type.startswith("image/"):
            fullText.append(extract_text_from_image(image.blob))
            
    return fullText

# Extract text in image file using pytesseract
def extract_text_from_image(img_bytes):
    # config for ocr engine and psm
    ocrConfig = r'--psm 6 --oem 3'
    image = PIL.Image.open(io.BytesIO(img_bytes))
    try:
        output = ocr.image_to_string(image, config=ocrConfig)
        return output
    finally:
        image.close()

# Extract text in pptx file using python-pptx
def extract_text_from_pptx(file):
    ppt = pptx.Presentation(io.BytesIO(file))
    fullText = []
    
    for slide in ppt.slides:
        for shape in slide.shapes:
            # Extract texts
            if hasattr(shape, 'text'):
                fullText.append(shape.text)
            # Extract images
            elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob  # the raw image data
                fullText.append(extract_text_from_image(image_bytes))
            # Extract Tables
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        fullText.append(cell.text)
    
    return fullText
